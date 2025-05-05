import sys
import PyPDF2
import docx
import pptx
import openpyxl
import os

def extract_text_from_pdf(pdf_path):
    text = []
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            if reader.is_encrypted:
                reader.decrypt("")  # Try to decrypt

            for page in reader.pages:
                extracted_text = page.extract_text()
                if extracted_text:
                    text.append(extracted_text)

        return "\n".join(text) if text else "Error: No extractable text (possibly an image-based PDF)"
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error extracting text from DOCX: {str(e)}"

def extract_text_from_pptx(pptx_path):
    try:
        prs = pptx.Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error extracting text from PPTX: {str(e)}"

def extract_text_from_xlsx(xlsx_path):
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text.append(" ".join([str(cell) for cell in row if cell is not None]))
        return "\n".join(text)
    except Exception as e:
        return f"Error extracting text from XLSX: {str(e)}"

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_path)
        elif ext == ".xlsx":
            return extract_text_from_xlsx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()
        else:
            return f"Unsupported file type: {ext}"
    except Exception as e:
        return f"Error extracting text: {str(e)}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Error: No file path provided.")
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' does not exist.")
        sys.exit(1)

    extracted_text = extract_text(file_path)
    
    # Force UTF-8 encoding for output
    sys.stdout.reconfigure(encoding="utf-8")
    print(extracted_text)
